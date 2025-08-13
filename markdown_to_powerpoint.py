#!/usr/bin/env python3
"""
Markdown to PowerPoint Converter

Converts bullet point markdown files to PowerPoint presentations.
Each ## header becomes a slide title with associated bullet points as slide content.
"""

import argparse
import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def parse_markdown_file(file_path):
    """Parse markdown file and extract headers with their bullet points."""
    sections = []
    current_section = None
    
    with open(file_path, 'r', encoding='utf-8') as file:
        for line in file:
            line = line.strip()
            
            # Skip empty lines and comments
            if not line or line.startswith('<!--'):
                continue
            
            # Check for header (## format)
            if line.startswith('## '):
                # Save previous section
                if current_section:
                    sections.append(current_section)
                
                # Start new section
                title = line[3:].strip()
                # Remove markdown bold formatting
                title = re.sub(r'\*\*(.*?)\*\*', r'\1', title)
                current_section = {
                    'title': title,
                    'bullets': []
                }
            
            # Check for bullet point
            elif line.startswith('• '):
                if current_section:
                    bullet_text = line[2:].strip()
                    current_section['bullets'].append(bullet_text)
        
        # Don't forget the last section
        if current_section:
            sections.append(current_section)
    
    return sections


def create_powerpoint(sections, output_file):
    """Create PowerPoint presentation from sections."""
    # Create presentation
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Generated Textbook Chapter Notes"
    subtitle.text = "Generated from Markdown Content"
    
    # Add content slides
    for section in sections:
        # Skip sections with no content
        if not section['bullets'] or (len(section['bullets']) == 1 and 
                                    'No content available' in section['bullets'][0]):
            continue
        
        # Use bullet layout
        bullet_slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section['title']
        
        # Add bullet points
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()  # Clear default text
        
        for i, bullet in enumerate(section['bullets']):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0  # Bullet level (0 = first level)
    
    # Save presentation
    prs.save(output_file)
    print(f"PowerPoint presentation saved as: {output_file}")


def main():
    """Main function to handle command line arguments and execute conversion."""
    parser = argparse.ArgumentParser(
        description='Convert bullet point markdown to PowerPoint presentation'
    )
    parser.add_argument(
        'input_file', 
        help='Input markdown file with bullet points'
    )
    parser.add_argument(
        '-o', '--output', 
        help='Output PowerPoint file name (default: input_file.pptx)',
        default=None
    )
    
    args = parser.parse_args()
    
    # Check if input file exists
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: Input file '{args.input_file}' not found.", file=sys.stderr)
        sys.exit(1)
    
    # Determine output filename
    if args.output:
        output_file = args.output
    else:
        output_file = input_path.stem + '.pptx'
    
    try:
        # Parse markdown file
        print(f"Parsing markdown file: {args.input_file}")
        sections = parse_markdown_file(input_path)
        print(f"Found {len(sections)} sections")
        
        # Create PowerPoint
        print(f"Creating PowerPoint presentation...")
        create_powerpoint(sections, output_file)
        print("Conversion completed successfully!")
        
    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()