from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from typing import Dict, Any, List, Optional

class PowerPointContentWriter:
    """
    Class for writing generated content to PowerPoint slides
    """
    
    def __init__(self):
        self.default_font_size = Pt(20)
        self.bullet_font_size = Pt(18)
    
    def add_content_to_slides(self, pptx_path: str, slides_with_content: List[Dict[str, Any]], 
                            output_path: Optional[str] = None) -> Dict[str, Any]:
        """
        Add generated content to PowerPoint slides
        
        Args:
            pptx_path: Path to original PowerPoint file
            slides_with_content: List of slides with generated content
            output_path: Output path for new PowerPoint file (if None, creates new filename)
        
        Returns:
            Dict with success status and details
        """
        
        try:
            # Load the presentation
            presentation = Presentation(pptx_path)
            
            # Determine output path
            if not output_path:
                base_name = os.path.splitext(pptx_path)[0]
                output_path = f"{base_name}_with_content.pptx"
            
            slides_processed = 0
            content_added = 0
            
            # Process each slide
            for slide_data in slides_with_content:
                slide_num = slide_data.get('slide_number', 0)
                
                # Skip if slide number is invalid
                if slide_num < 1 or slide_num > len(presentation.slides):
                    continue
                
                slide = presentation.slides[slide_num - 1]  # Convert to 0-based index
                generated_content = slide_data.get('generated_content')
                
                if generated_content and generated_content.get('success'):
                    success = self._add_content_to_slide(slide, generated_content)
                    if success:
                        content_added += 1
                
                slides_processed += 1
            
            # Save the presentation
            presentation.save(output_path)
            
            return {
                'success': True,
                'output_path': output_path,
                'slides_processed': slides_processed,
                'content_added': content_added,
                'message': f"Successfully added content to {content_added} out of {slides_processed} slides"
            }
        
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f"Failed to add content to PowerPoint: {e}"
            }
    
    def _add_content_to_slide(self, slide, generated_content: Dict[str, Any]) -> bool:
        """
        Add generated content to a specific slide
        
        Args:
            slide: PowerPoint slide object
            generated_content: Generated content dictionary
        
        Returns:
            True if content was added successfully
        """
        
        try:
            bullet_points = generated_content.get('bullet_points', [])
            topic = generated_content.get('topic', '')
            
            if not bullet_points:
                return False
            
            # Find the best place to add content
            content_placeholder = self._find_content_placeholder(slide)
            
            if content_placeholder:
                # Add to existing content placeholder
                self._add_bullets_to_placeholder(content_placeholder, bullet_points)
                return True
            else:
                # Create new text box for content
                return self._create_content_textbox(slide, bullet_points, topic)
        
        except Exception as e:
            print(f"Error adding content to slide: {e}")
            return False
    
    def _find_content_placeholder(self, slide):
        """Find the best placeholder for adding content"""
        
        # Look for content placeholders
        for shape in slide.shapes:
            if shape.is_placeholder:
                # Common content placeholder types
                if shape.placeholder_format.type in [1, 2, 8, 14]:  # Content, text, object, content with caption
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        return shape
        
        # Look for existing text boxes that might be suitable
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                # If text frame is mostly empty, we can use it
                if len(shape.text_frame.text.strip()) < 50:
                    return shape
        
        return None
    
    def _add_bullets_to_placeholder(self, placeholder, bullet_points: List[str]):
        """Add bullet points to an existing placeholder"""
        
        text_frame = placeholder.text_frame
        text_frame.clear()  # Clear existing content
        
        # Set up the text frame
        text_frame.word_wrap = True
        text_frame.auto_size = True
        
        for i, bullet in enumerate(bullet_points):
            if i == 0:
                # Use the existing paragraph for the first bullet
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraph for subsequent bullets
                p = text_frame.add_paragraph()
            
            # Set bullet point
            p.text = bullet
            p.level = 0  # Top level bullet
            
            # Format the paragraph
            if p.runs:
                run = p.runs[0]
                run.font.size = self.bullet_font_size
                run.font.name = 'Arial'
    
    def _create_content_textbox(self, slide, bullet_points: List[str], topic: str) -> bool:
        """Create a new text box for content"""
        
        try:
            # Calculate position - try to place in a good spot
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(4)
            
            # Check if this position conflicts with existing shapes
            if self._position_conflicts_with_existing(slide, left, top, width, height):
                # Try alternative positions
                positions = [
                    (Inches(0.5), Inches(1.5), Inches(9), Inches(5)),  # Wider, higher
                    (Inches(1.5), Inches(3), Inches(7), Inches(3.5)),  # Narrower, lower
                    (Inches(0.5), Inches(4), Inches(9), Inches(3)),    # Bottom area
                ]
                
                for pos in positions:
                    if not self._position_conflicts_with_existing(slide, *pos):
                        left, top, width, height = pos
                        break
            
            # Create text box
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Add bullet points
            for i, bullet in enumerate(bullet_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
                
                # Format
                if p.runs:
                    run = p.runs[0]
                    run.font.size = self.bullet_font_size
                    run.font.name = 'Arial'
                    run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            
            return True
        
        except Exception as e:
            print(f"Error creating content textbox: {e}")
            return False
    
    def _position_conflicts_with_existing(self, slide, left, top, width, height) -> bool:
        """Check if a position conflicts with existing shapes"""
        
        new_right = left + width
        new_bottom = top + height
        
        for shape in slide.shapes:
            # Skip title shapes and very small shapes
            if (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and 
                hasattr(shape, 'placeholder_format') and 
                shape.placeholder_format.type == 0):  # Title placeholder
                continue
            
            if shape.width < Inches(0.5) or shape.height < Inches(0.5):
                continue
            
            # Check for overlap
            shape_right = shape.left + shape.width
            shape_bottom = shape.top + shape.height
            
            if not (new_right <= shape.left or left >= shape_right or 
                   new_bottom <= shape.top or top >= shape_bottom):
                return True  # Conflict detected
        
        return False
    
    def create_content_summary_slide(self, presentation_path: str, content_results: Dict[str, Any], 
                                   output_path: Optional[str] = None) -> Dict[str, Any]:
        """
        Create a summary slide showing what content was generated
        
        Args:
            presentation_path: Path to PowerPoint file
            content_results: Results from content generation
            output_path: Output path for new PowerPoint file
        
        Returns:
            Dict with success status and details
        """
        
        try:
            presentation = Presentation(presentation_path)
            
            # Add new slide at the end
            slide_layout = presentation.slide_layouts[1]  # Title and content layout
            slide = presentation.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = "Generated Slide Content Summary"
            
            # Add summary content
            content = slide.placeholders[1]  # Content placeholder
            text_frame = content.text_frame
            text_frame.clear()
            
            # Add summary information
            summary_lines = [
                f"Total slides processed: {content_results.get('total_slides', 0)}",
                f"Content successfully generated: {content_results.get('successful', 0)}",
                f"Generation method: {content_results.get('generation_method', 'unknown')}",
                f"Success rate: {content_results.get('success_rate', 0):.1%}",
                "",
                "Generated content follows 6x6 rule:",
                "• Maximum 6 words per bullet point",
                "• Maximum 6 bullet points per topic",
                "",
                "Content based on:",
                "• Slide titles and headings",
                "• Aligned textbook sections",
                "• Educational best practices"
            ]
            
            for i, line in enumerate(summary_lines):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = line
                
                if line.startswith("•"):
                    p.level = 1
                else:
                    p.level = 0
                
                # Format
                if p.runs:
                    run = p.runs[0]
                    run.font.size = Pt(16)
                    run.font.name = 'Arial'
            
            # Save presentation
            if not output_path:
                base_name = os.path.splitext(presentation_path)[0]
                output_path = f"{base_name}_with_summary.pptx"
            
            presentation.save(output_path)
            
            return {
                'success': True,
                'output_path': output_path,
                'message': f"Summary slide added successfully"
            }
        
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'message': f"Failed to create summary slide: {e}"
            }