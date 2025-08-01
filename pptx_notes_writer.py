#!/usr/bin/env python3

import os
import re
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches

class PowerPointNotesWriter:
    """
    Writes generated notes back to PowerPoint slide notes sections
    """
    
    def __init__(self):
        self.supported_extensions = ['.pptx', '.ppt']
    
    def write_notes_to_pptx(self, pptx_path: str, notes_content: str, 
                           output_path: Optional[str] = None) -> Dict[str, Any]:
        """
        Write generated notes to PowerPoint slide notes sections
        
        Args:
            pptx_path: Path to original PowerPoint file
            notes_content: Generated notes content (markdown format)
            output_path: Path for output file (if None, overwrites original)
        
        Returns:
            Dict with success status and details
        """
        
        try:
            # Load the presentation
            if not os.path.exists(pptx_path):
                return {
                    'success': False,
                    'error': f'PowerPoint file not found: {pptx_path}'
                }
            
            prs = Presentation(pptx_path)
            
            # Parse notes content to extract slide-specific notes
            slide_notes = self._parse_notes_by_slide(notes_content)
            
            # Write notes to each slide
            slides_updated = 0
            for slide_idx, slide in enumerate(prs.slides):
                slide_number = slide_idx + 1
                
                if str(slide_number) in slide_notes:
                    # Get the notes slide
                    notes_slide = slide.notes_slide
                    
                    # Clear existing notes
                    notes_slide.notes_text_frame.clear()
                    
                    # Add new notes
                    slide_note_content = slide_notes[str(slide_number)]
                    notes_slide.notes_text_frame.text = slide_note_content
                    
                    slides_updated += 1
            
            # Save the presentation
            save_path = output_path or pptx_path
            prs.save(save_path)
            
            return {
                'success': True,
                'slides_updated': slides_updated,
                'total_slides': len(prs.slides),
                'output_path': save_path,
                'message': f'Successfully updated notes for {slides_updated} slides'
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'error_type': type(e).__name__
            }
    
    def _parse_notes_by_slide(self, notes_content: str) -> Dict[str, str]:
        """
        Parse markdown notes content and extract notes for each slide
        
        Args:
            notes_content: Full notes content in markdown format
        
        Returns:
            Dict mapping slide numbers to their notes content
        """
        
        slide_notes = {}
        current_slide = None
        current_notes = []
        
        lines = notes_content.split('\n')
        
        for line in lines:
            # Check for slide headers
            slide_match = re.match(r'## Slide (\d+):', line)
            if slide_match:
                # Save previous slide notes if any
                if current_slide and current_notes:
                    slide_notes[current_slide] = self._format_notes_for_pptx(current_notes)
                
                # Start new slide
                current_slide = slide_match.group(1)
                current_notes = []
                # Clean title: remove "## " and "Slide #: " prefix
                clean_title = line.replace('## ', '').replace('#', '')
                # Remove "Slide X: " pattern
                clean_title = re.sub(r'^Slide \d+:\s*', '', clean_title)
                current_notes.append(clean_title)
                current_notes.append('')  # Add line feed after title
                continue
            
            # Skip metadata, table of contents, and summary sections
            if line.startswith('---') or line.startswith('# Presentation Notes') or \
               line.startswith('## Table of Contents') or line.startswith('## Presentation Summary') or \
               line.startswith('**Generated on:**') or line.startswith('- [Slide'):
                continue
            
            # Skip HTML comments
            if line.strip().startswith('<!--') and line.strip().endswith('-->'):
                continue
            
            # Collect slide content
            if current_slide and line.strip():
                current_notes.append(line)
        
        # Don't forget the last slide
        if current_slide and current_notes:
            slide_notes[current_slide] = self._format_notes_for_pptx(current_notes)
        
        return slide_notes
    
    def _format_notes_for_pptx(self, notes_lines: List[str]) -> str:
        """
        Format notes lines for PowerPoint notes section
        
        Args:
            notes_lines: List of note lines for a slide
        
        Returns:
            Formatted notes text
        """
        
        formatted_lines = []
        
        for line in notes_lines:
            # Handle Grok format: **I. content** -> I. content
            if line.strip().startswith('**I.') or line.strip().startswith('**II.') or \
               line.strip().startswith('**III.') or line.strip().startswith('**IV.') or \
               line.strip().startswith('**V.') or line.strip().startswith('**VI.'):
                # Remove ** formatting and keep the content
                clean_line = line.replace('**', '')
                formatted_lines.append(clean_line)
                continue
            
            # Handle sub-points: "  A. content" or "  B. content" - keep exactly as is
            if line.strip().startswith('A. ') or line.strip().startswith('B. '):
                # Ensure exactly 2 spaces for indentation
                content = line.strip()
                formatted_lines.append('  ' + content)
                
                # Add extra line break after B sub-points for better spacing
                if content.startswith('B. '):
                    formatted_lines.append('')  # Empty line for spacing
                continue
            
            # Remove other markdown formatting
            clean_line = line.replace('**', '').replace('*', '')
            
            # Add all lines (including empty ones for spacing)
            formatted_lines.append(clean_line)
        
        return '\n'.join(formatted_lines)
    
    def create_notes_only_version(self, pptx_path: str, notes_content: str, 
                                 output_path: str) -> Dict[str, Any]:
        """
        Create a new PowerPoint file with only the notes (no slide content changes)
        
        Args:
            pptx_path: Path to original PowerPoint file
            notes_content: Generated notes content
            output_path: Path for new PowerPoint file with notes
        
        Returns:
            Dict with success status and details
        """
        
        try:
            # Create a copy with notes
            result = self.write_notes_to_pptx(pptx_path, notes_content, output_path)
            
            if result['success']:
                result['message'] = f"Created new PowerPoint file with notes: {output_path}"
            
            return result
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'error_type': type(e).__name__
            }
    
    def extract_existing_notes(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract existing notes from a PowerPoint file
        
        Args:
            pptx_path: Path to PowerPoint file
        
        Returns:
            Dict with existing notes by slide
        """
        
        try:
            prs = Presentation(pptx_path)
            existing_notes = {}
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_number = slide_idx + 1
                notes_text = slide.notes_slide.notes_text_frame.text
                
                if notes_text.strip():
                    existing_notes[str(slide_number)] = notes_text
            
            return {
                'success': True,
                'notes': existing_notes,
                'total_slides': len(prs.slides),
                'slides_with_notes': len(existing_notes)
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'error_type': type(e).__name__
            }