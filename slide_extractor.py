from pptx import Presentation
from typing import List, Dict, Any
import re

class SlideExtractor:
    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.presentation = None
        
    def load_presentation(self) -> bool:
        try:
            self.presentation = Presentation(self.pptx_path)
            return True
        except Exception as e:
            print(f"Error loading presentation: {e}")
            return False
    
    def extract_slide_content(self, slide_index: int) -> Dict[str, Any]:
        if not self.presentation or slide_index >= len(self.presentation.slides):
            return {}
            
        slide = self.presentation.slides[slide_index]
        content = {
            'slide_number': slide_index + 1,
            'title': '',
            'text_content': [],
            'bullet_points': [],
            'images': [],
            'notes': ''
        }
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                
                if shape == slide.shapes.title:
                    content['title'] = text
                else:
                    content['text_content'].append(text)
                    
                    if self._is_bullet_point(text):
                        content['bullet_points'].extend(self._extract_bullets(text))
            
            if shape.shape_type == 13:  # Picture
                content['images'].append({
                    'name': shape.name,
                    'width': shape.width,
                    'height': shape.height
                })
        
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame.text:
                content['notes'] = notes_slide.notes_text_frame.text.strip()
        
        return content
    
    def extract_all_slides(self) -> List[Dict[str, Any]]:
        if not self.load_presentation():
            return []
            
        slides_content = []
        for i in range(len(self.presentation.slides)):
            slide_content = self.extract_slide_content(i)
            if slide_content:
                slides_content.append(slide_content)
                
        return slides_content
    
    def _is_bullet_point(self, text: str) -> bool:
        bullet_patterns = [r'^\s*[-•*]\s', r'^\s*\d+\.\s', r'^\s*[a-zA-Z]\.\s']
        return any(re.match(pattern, line) for line in text.split('\n') for pattern in bullet_patterns)
    
    def _extract_bullets(self, text: str) -> List[str]:
        bullets = []
        lines = text.split('\n')
        for line in lines:
            line = line.strip()
            if line and re.match(r'^\s*[-•*]\s|^\s*\d+\.\s|^\s*[a-zA-Z]\.\s', line):
                clean_bullet = re.sub(r'^\s*[-•*]\s*|\s*\d+\.\s*|\s*[a-zA-Z]\.\s*', '', line)
                bullets.append(clean_bullet)
        return bullets
    
    def get_slide_summary(self, slide_index: int) -> str:
        content = self.extract_slide_content(slide_index)
        if not content:
            return ""
            
        summary_parts = []
        if content['title']:
            summary_parts.append(f"Title: {content['title']}")
        
        if content['bullet_points']:
            summary_parts.append(f"Key Points: {'; '.join(content['bullet_points'][:3])}")
        
        if content['text_content']:
            all_text = ' '.join(content['text_content'])
            if len(all_text) > 200:
                all_text = all_text[:200] + "..."
            summary_parts.append(f"Content: {all_text}")
            
        return " | ".join(summary_parts)